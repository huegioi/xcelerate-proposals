#!/usr/bin/env python3
"""
Xcelerate Growth Partners - Proposal Deck Generator

Produces TWO files for every proposal:
  1. A PDF  (same filename, .pdf extension)
  2. A PPTX (same filename, .pptx extension)

Usage:
    python generate_proposal.py \
        --company   "Thornburg Investments" \
        --contact   "Sarah Chen" \
        --date      "September 15, 2025" \
        --base-pdf  "/path/to/base_proposal_template.pdf" \
        --services  "Team Development and Optimization|Succession Planning|Keynote Speeches" \
        --costs     "Program Fee: $25,000|Travel & Expenses: At cost" \
        --notes     "Optional custom paragraph for the services slide." \
        --output    "/path/to/output/Thornburg_Proposal.pdf"

Services and costs use | as a separator between line items.
The .pptx file is saved automatically next to the .pdf output.
"""

import argparse
import os
import io

# ── reportlab (PDF) ───────────────────────────────────────────────────────────
from reportlab.pdfgen import canvas
from reportlab.lib.colors import HexColor, white
from reportlab.lib.utils import ImageReader
from reportlab.platypus import Paragraph
from reportlab.lib.styles import ParagraphStyle
from pypdf import PdfReader, PdfWriter

# ── python-pptx (PPTX) ───────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# ── Shared brand constants ────────────────────────────────────────────────────
# ReportLab colours
RL_NAVY      = HexColor("#1D3461")
RL_GREEN     = HexColor("#5CB85C")
RL_GREEN_DK  = HexColor("#3A7A3A")
RL_WHITE     = white
RL_BODY      = HexColor("#2C2C2C")
RL_LGRAY     = HexColor("#F5F5F5")
RL_MGRAY     = HexColor("#CCCCCC")

# python-pptx colours
PT_NAVY      = RGBColor(0x1D, 0x34, 0x61)
PT_GREEN     = RGBColor(0x5C, 0xB8, 0x5C)
PT_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
PT_BODY      = RGBColor(0x2C, 0x2C, 0x2C)
PT_LGRAY     = RGBColor(0xF5, 0xF5, 0xF5)
PT_MGRAY     = RGBColor(0xCC, 0xCC, 0xCC)

# PDF page size — matches the base Edward Jones proposal PDF exactly
PDF_W, PDF_H = 720.0, 405.0

# PPTX slide size — standard 16:9 widescreen (13.33" x 7.5")
PPTX_W = Inches(13.33)
PPTX_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════════
#  Helpers
# ═══════════════════════════════════════════════════════════════════════════════

def _logo_path(script_dir: str) -> str:
    return os.environ.get(
        "XCELERATE_LOGO",
        os.path.join(script_dir, "..", "assets", "xcelerate_logo.png")
    )


def _pptx_output_path(pdf_output: str) -> str:
    """Derive the .pptx path from the .pdf output path."""
    base = os.path.splitext(pdf_output)[0]
    return base + ".pptx"


# ── Boilerplate slide content (mirrors the Edward Jones proposal template) ────
STANDARD_SLIDES = [
    {
        "title": "Who We Are",
        "bullets": [
            "We are a team with extensive knowledge in Wealth Management. Our expertise is broad and diverse, with a unique understanding of clients, advisors and investment products.",
            "We are keenly focused on growth and business development. Each engagement is custom, and we partner with clients to drive specific desired outcomes.",
            "Our capabilities are unique and differentiated. We merge the talents of our core team with a roster of strategic alliances to ensure best-in-class support and project execution.",
            "Our process and strategy is to excel at leading edge thinking, converting fresh new ideas into measurable results and positive transformation.",
        ],
    },
    {
        "title": "Our Competitive Advantage",
        "bullets": [
            "Team Development and Optimization",
            "Succession Planning and Practice Transitions",
            "The Future of Wealth Management",
            "Developing Leadership",
            "Support Staff Development and Career Pathing",
            "The Wealth Management Process",
            "  – Client Management",
            "  – Financial Planning",
            "  – The Investment Process",
            "  – Exceeding Client Expectations",
        ],
    },
    {
        "title": "Why Xcelerate?",
        "subtitle": "Asset Manager Opportunities",
        "bullets": [
            "Over 150 years of significant Wealth Management experience.",
            "Team's track record of success and industry credibility.",
            "Access to cutting edge and relevant content.",
            "Firsthand knowledge of Financial Advisor and client preferences.",
            "Access to talented presenters and thought-provoking content development specialists.",
            "Compelling conversations about The Future of Wealth Management.",
            "Creative strategies to improve engagement and outcomes.",
        ],
    },
    {
        "title": "Specific Partnership Ideas",
        "bullets": [
            "Working with Wealth Management firms on strategy and growth around specific themes such as retirement planning, succession planning or acting on key business initiatives.",
            "Content development and value-added white papers.",
            "Interactive and innovative idea development.",
            "Cutting edge Practice Management training by talented industry professionals.",
            "Keynote speeches.",
            "Conference and sponsorship strategy to enhance relationships and outcomes.",
            "Action plan to align annual spend and year-long marketing planning.",
            "Advisor relationship development.",
            '"Train the trainer" workshops to improve wholesaler presentations.',
            "Access to best practice thinking, curated articles and creative training initiatives.",
            "Production of webinars, podcasts, and access to sought-after programs.",
        ],
    },
    {
        "title": "Our Team of Experts",
        "team": [
            ("Jim Tracy", "Co-Founder & CEO",
             "As a 40-year veteran of the wealth management industry, Jim is a recognized leader. He has achieved considerable success in growing and building innovative solutions that benefit advisors and their clients."),
            ("Mary Deatherage", "Co-Founder & President",
             "Mary is a respected industry leader, who previously led a multi-billion dollar team at Morgan Stanley. Her strategic vision and expertise are invaluable assets to the firm."),
            ("Tara Forrest", "Head of Business Management",
             "Tara brings a wealth of expertise in developing and implementing effective business strategies. Her background in operations and risk management add structure and additional capabilities to the group."),
        ],
    },
]


# ═══════════════════════════════════════════════════════════════════════════════
#  PDF GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

def _pdf_footer(c, page_num):
    c.setFillColor(RL_NAVY)
    c.setFont("Helvetica-Bold", 9)
    c.drawString(40, 18, "Xcelerate Growth Partners")
    c.setFillColor(RL_GREEN)
    c.rect(PDF_W - 50, 8, 36, 24, fill=1, stroke=0)
    c.setFillColor(RL_WHITE)
    c.setFont("Helvetica-Bold", 11)
    c.drawCentredString(PDF_W - 32, 16, str(page_num))


def _pdf_cover(c, company, contact, date, logo_path):
    mid = PDF_W / 2
    c.setFillColor(RL_NAVY)
    c.rect(0, 0, PDF_W, PDF_H, fill=1, stroke=0)

    # ── Logo (top band) ───────────────────────────────────────────────────────
    if logo_path and os.path.exists(logo_path):
        c.drawImage(ImageReader(logo_path),
                    mid - 100, PDF_H - 125, 200, 85,
                    mask='auto', preserveAspectRatio=True, anchor='c')

    # ── Green accent bar below logo ───────────────────────────────────────────
    c.setFillColor(RL_GREEN)
    c.rect(0, PDF_H - 135, PDF_W, 6, fill=1, stroke=0)

    # ── Contact / Company name (well below logo) ──────────────────────────────
    c.setFillColor(RL_WHITE)
    if contact:
        c.setFont("Helvetica-Bold", 34)
        c.drawCentredString(mid, 225, contact)
        c.setFont("Helvetica-Bold", 28)
        c.drawCentredString(mid, 188, company)
    else:
        c.setFont("Helvetica-Bold", 34)
        c.drawCentredString(mid, 210, company)

    # ── Delivery Launch Date label + date ─────────────────────────────────────
    c.setFillColor(RL_WHITE)
    c.setFont("Helvetica", 9)
    c.drawCentredString(mid, 158, "Delivery Launch Date:")
    c.setFillColor(RL_GREEN)
    c.setFont("Helvetica-Bold", 20)
    c.drawCentredString(mid, 140, date)

    # ── Presented By — centred ────────────────────────────────────────────────
    c.setFillColor(RL_WHITE)
    c.setFont("Helvetica", 9)
    c.drawCentredString(mid, 100, "Presented By:")
    c.setFont("Helvetica-Bold", 9)
    c.drawCentredString(mid, 84, "Jim Tracy, Co-Founder & CEO")
    c.drawCentredString(mid, 68, "Mary Deatherage, Co-Founder & President")
    c.setFont("Helvetica", 9)
    c.drawCentredString(mid, 52, "Xcelerate Growth Partners")

    # ── Bottom green bar ──────────────────────────────────────────────────────
    c.setFillColor(RL_GREEN)
    c.rect(0, 0, PDF_W, 8, fill=1, stroke=0)
    c.showPage()


def _pdf_services(c, company, services, notes, page_num):
    c.setFillColor(RL_WHITE)
    c.rect(0, 0, PDF_W, PDF_H, fill=1, stroke=0)
    c.setFillColor(RL_NAVY)
    c.setFont("Helvetica-Bold", 20)
    c.drawString(36, PDF_H - 44, "Recommended Services")
    c.setFillColor(RL_BODY)
    c.setFont("Helvetica", 10)
    c.drawString(36, PDF_H - 60, f"Tailored engagement plan for {company}")
    c.setFillColor(RL_GREEN)
    c.rect(36, PDF_H - 68, PDF_W - 72, 2, fill=1, stroke=0)
    bullet_style = ParagraphStyle('b', fontName='Helvetica-Bold', fontSize=10,
                                  leading=14, textColor=RL_NAVY, leftIndent=12)
    mid = PDF_W / 2
    col_w = mid - 60
    y_l = y_r = PDF_H - 84
    for i, svc in enumerate(services):
        p = Paragraph(f"• {svc}", bullet_style)
        pw, ph = p.wrap(col_w, PDF_H)
        if i < (len(services) + 1) // 2:
            p.drawOn(c, 36, y_l - ph);  y_l -= ph + 6
        else:
            p.drawOn(c, mid + 10, y_r - ph);  y_r -= ph + 6
    if notes:
        ns = ParagraphStyle('n', fontName='Helvetica-Oblique', fontSize=9,
                            leading=13, textColor=RL_BODY)
        p = Paragraph(notes, ns)
        pw, ph = p.wrap(PDF_W - 72, PDF_H)
        p.drawOn(c, 36, min(y_l, y_r) - 10 - ph)
    _pdf_footer(c, page_num)
    c.showPage()


def _pdf_investment(c, company, costs, page_num):
    c.setFillColor(RL_WHITE)
    c.rect(0, 0, PDF_W, PDF_H, fill=1, stroke=0)
    c.setFillColor(RL_NAVY)
    c.setFont("Helvetica-Bold", 20)
    c.drawString(36, PDF_H - 44, "Proposed Investment")
    c.setFillColor(RL_BODY)
    c.setFont("Helvetica", 10)
    c.drawString(36, PDF_H - 60, f"Fee structure for {company}")
    c.setFillColor(RL_GREEN)
    c.rect(36, PDF_H - 68, PDF_W - 72, 2, fill=1, stroke=0)
    row_h, row_y = 30, PDF_H - 84
    # Separate out summary/total lines from line items
    summary_keys = ("total investment", "monthly retainer", "one-time fees")
    item_costs   = [l for l in costs if not any(l.lower().startswith(k) for k in summary_keys)]
    summary_lines = [l for l in costs if any(l.lower().startswith(k) for k in summary_keys)]
    total_lines   = [l for l in costs if l.lower().startswith("total")]
    subtotal_lines= [l for l in summary_lines if not l.lower().startswith("total")]

    for i, line in enumerate(item_costs):
        label, amount = (line.split(":", 1) if ":" in line else (line, ""))
        if i % 2 == 0:
            c.setFillColor(RL_LGRAY)
            c.rect(36, row_y - row_h + 8, PDF_W - 72, row_h, fill=1, stroke=0)
        c.setFillColor(RL_NAVY);  c.setFont("Helvetica-Bold", 11)
        c.drawString(48, row_y - 6, label.strip())
        c.setFillColor(RL_GREEN_DK); c.setFont("Helvetica-Bold", 11)
        c.drawRightString(PDF_W - 48, row_y - 6, amount.strip())
        row_y -= row_h
    # Subtotal rows (Monthly Retainer / One-Time) with slight indent
    if subtotal_lines:
        row_y -= 4
        for line in subtotal_lines:
            label, amount = (line.split(":", 1) if ":" in line else (line, ""))
            c.setFillColor(RL_LGRAY)
            c.rect(36, row_y - row_h + 10, PDF_W - 72, row_h - 4, fill=1, stroke=0)
            c.setFillColor(RL_NAVY);  c.setFont("Helvetica-BoldOblique", 10)
            c.drawString(52, row_y - 5, label.strip())
            c.setFillColor(RL_GREEN_DK); c.setFont("Helvetica-Bold", 10)
            c.drawRightString(PDF_W - 48, row_y - 5, amount.strip())
            row_y -= (row_h - 4)
    # Total row — navy background, larger text
    if total_lines:
        label, amount = (total_lines[0].split(":", 1) if ":" in total_lines[0] else (total_lines[0], ""))
        row_y -= 6
        c.setFillColor(RL_NAVY)
        c.rect(36, row_y - row_h + 8, PDF_W - 72, row_h + 4, fill=1, stroke=0)
        c.setFillColor(RL_WHITE); c.setFont("Helvetica-Bold", 13)
        c.drawString(48, row_y - 5, label.strip())
        c.setFillColor(RL_GREEN); c.setFont("Helvetica-Bold", 13)
        c.drawRightString(PDF_W - 48, row_y - 5, amount.strip())
    c.setFillColor(RL_MGRAY);  c.setFont("Helvetica-Oblique", 8)
    c.drawString(36, 48, "All fees are subject to final scope confirmation. "
                 "Travel expenses billed at cost.")
    _pdf_footer(c, page_num)
    c.showPage()


def build_proposal_pdf(company, contact, date, base_pdf_path,
                       services, costs, notes, output_path, logo_path):
    writer = PdfWriter()

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=(PDF_W, PDF_H))
    _pdf_cover(c, company, contact, date, logo_path)
    c.save();  buf.seek(0)
    writer.add_page(PdfReader(buf).pages[0])

    base_count = 0
    if base_pdf_path and os.path.exists(base_pdf_path):
        base_reader = PdfReader(base_pdf_path)
        for page in base_reader.pages[1:]:
            writer.add_page(page)
        base_count = len(base_reader.pages) - 1
    else:
        print("Warning: base PDF not found — skipping template pages")

    if services:
        buf2 = io.BytesIO()
        c2 = canvas.Canvas(buf2, pagesize=(PDF_W, PDF_H))
        _pdf_services(c2, company, services, notes, base_count + 2)
        c2.save();  buf2.seek(0)
        writer.add_page(PdfReader(buf2).pages[0])

    if costs:
        buf3 = io.BytesIO()
        c3 = canvas.Canvas(buf3, pagesize=(PDF_W, PDF_H))
        _pdf_investment(c3, company, costs, base_count + 3)
        c3.save();  buf3.seek(0)
        writer.add_page(PdfReader(buf3).pages[0])

    with open(output_path, "wb") as f:
        writer.write(f)
    print(f"PDF  saved → {output_path}")


# ═══════════════════════════════════════════════════════════════════════════════
#  PPTX GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

# ── Content-slide layout constants ───────────────────────────────────────────
_HDR_H   = Inches(0.88)   # navy header band height
_GRN_H   = Inches(0.06)   # green accent bar below header
_CTT_Y   = _HDR_H + _GRN_H   # where content area starts (Inches ~0.94)


def _pptx_blank_slide(prs):
    """Add a completely blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]   # index 6 = Blank
    return prs.slides.add_slide(blank_layout)


def _pptx_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _pptx_rect(slide, left, top, width, height, fill_color: RGBColor,
               line_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _pptx_textbox(slide, left, top, width, height, text, font_size,
                  bold=False, color: RGBColor = None, align=PP_ALIGN.LEFT,
                  italic=False, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb


def _pptx_footer(slide, page_num):
    W, H = PPTX_W, PPTX_H
    # "Xcelerate Growth Partners" in GREEN (matches PDF template style)
    _pptx_textbox(slide, W - Inches(5.2), H - Inches(0.58),
                  Inches(4.5), Inches(0.50),
                  "Xcelerate Growth Partners", 18,
                  bold=True, color=PT_GREEN, align=PP_ALIGN.RIGHT)
    # Green page-number box
    _pptx_rect(slide,
               W - Inches(0.72), H - Inches(0.58),
               Inches(0.58), Inches(0.50), PT_GREEN)
    pn_box = _pptx_textbox(slide, W - Inches(0.72), H - Inches(0.58),
                           Inches(0.58), Inches(0.50),
                           str(page_num), 18,
                           bold=True, color=PT_WHITE, align=PP_ALIGN.CENTER)
    pn_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _pptx_slide_header(slide, title, logo_path, subtitle=None):
    """
    Draws the branded navy header on every content slide:
      - Full-width navy band
      - Title in white on the left
      - Logo on the right
      - Thin green accent bar underneath
      - Optional subtitle in dark text just below the green bar
    """
    W = PPTX_W
    _pptx_rect(slide, 0, 0, W, _HDR_H, PT_NAVY)

    # Title text (white, left-aligned inside header)
    _pptx_textbox(slide, Inches(0.4), Inches(0.14),
                  W - Inches(2.6), Inches(0.62),
                  title, 22, bold=True, color=PT_WHITE)

    # Logo top-right of header
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path,
            W - Inches(2.25), Inches(0.09),
            Inches(2.0), Inches(0.7))

    # Green accent bar below header
    _pptx_rect(slide, 0, _HDR_H, W, _GRN_H, PT_GREEN)

    # Optional subtitle just below the green bar
    if subtitle:
        _pptx_textbox(slide, Inches(0.4), _CTT_Y + Inches(0.06),
                      W - Inches(1), Inches(0.3),
                      subtitle, 11, color=PT_BODY)


def _pptx_cover(prs, company, contact, date, logo_path):
    slide = _pptx_blank_slide(prs)
    W, H = PPTX_W, PPTX_H
    _pptx_bg(slide, PT_NAVY)

    # ── Logo (top, centred — preserve aspect ratio, never stretch) ───────────
    if logo_path and os.path.exists(logo_path):
        # Only specify width; python-pptx auto-calculates height to preserve ratio
        logo_w = Inches(2.2)
        pic = slide.shapes.add_picture(logo_path,
                                       W / 2 - logo_w / 2, Inches(0.20),
                                       width=logo_w)

    # ── Green accent bar below logo ───────────────────────────────────────────
    _pptx_rect(slide, 0, Inches(1.72), W, Inches(0.07), PT_GREEN)

    # ── Contact / Company name ────────────────────────────────────────────────
    if contact:
        _pptx_textbox(slide, Inches(0.75), Inches(1.92),
                      W - Inches(1.5), Inches(0.85),
                      contact, 48,
                      bold=True, color=PT_WHITE, align=PP_ALIGN.CENTER)
        _pptx_textbox(slide, Inches(0.75), Inches(2.85),
                      W - Inches(1.5), Inches(0.75),
                      company, 40,
                      bold=True, color=PT_WHITE, align=PP_ALIGN.CENTER)
    else:
        _pptx_textbox(slide, Inches(0.75), Inches(2.2),
                      W - Inches(1.5), Inches(0.85),
                      company, 48,
                      bold=True, color=PT_WHITE, align=PP_ALIGN.CENTER)

    # ── Delivery Launch Date label + date ─────────────────────────────────────
    date_label_y = Inches(3.78) if contact else Inches(3.35)
    _pptx_textbox(slide, Inches(0.75), date_label_y,
                  W - Inches(1.5), Inches(0.40),
                  "Delivery Launch Date:", 14,
                  bold=False, color=PT_WHITE, align=PP_ALIGN.CENTER)
    _pptx_textbox(slide, Inches(0.75), date_label_y + Inches(0.42),
                  W - Inches(1.5), Inches(0.65),
                  date, 30,
                  bold=True, color=PT_GREEN, align=PP_ALIGN.CENTER)

    # ── Presented By — centred ────────────────────────────────────────────────
    by_top = H - Inches(1.90)
    _pptx_textbox(slide, Inches(0.75), by_top,
                  W - Inches(1.5), Inches(0.38),
                  "Presented By:", 14,
                  bold=False, color=PT_WHITE, align=PP_ALIGN.CENTER)
    _pptx_textbox(slide, Inches(0.75), by_top + Inches(0.40),
                  W - Inches(1.5), Inches(0.38),
                  "Jim Tracy, Co-Founder & CEO", 16,
                  bold=True, color=PT_WHITE, align=PP_ALIGN.CENTER)
    _pptx_textbox(slide, Inches(0.75), by_top + Inches(0.80),
                  W - Inches(1.5), Inches(0.38),
                  "Mary Deatherage, Co-Founder & President", 16,
                  bold=True, color=PT_WHITE, align=PP_ALIGN.CENTER)
    _pptx_textbox(slide, Inches(0.75), by_top + Inches(1.20),
                  W - Inches(1.5), Inches(0.38),
                  "Xcelerate Growth Partners", 14,
                  bold=False, color=PT_WHITE, align=PP_ALIGN.CENTER)

    # ── Bottom green bar ──────────────────────────────────────────────────────
    _pptx_rect(slide, 0, H - Inches(0.15), W, Inches(0.15), PT_GREEN)


def _pptx_content_slide(prs, title, bullets, subtitle=None,
                        page_num=None, two_col=False, logo_path=None):
    slide = _pptx_blank_slide(prs)
    W, H = PPTX_W, PPTX_H
    _pptx_bg(slide, PT_WHITE)

    # ── Large navy title (matches PDF template) ───────────────────────────────
    _pptx_textbox(slide, Inches(0.45), Inches(0.22),
                  W - Inches(0.9), Inches(0.92),
                  title, 32, bold=True, color=PT_NAVY)

    # Optional subtitle just below title
    title_bottom = Inches(1.18)
    if subtitle:
        _pptx_textbox(slide, Inches(0.45), Inches(1.18),
                      W - Inches(0.9), Inches(0.36),
                      subtitle, 14, color=PT_BODY)
        title_bottom = Inches(1.56)

    # Thin green divider
    _pptx_rect(slide, Inches(0.45), title_bottom,
               W - Inches(0.9), Inches(0.04), PT_GREEN)

    y_start = title_bottom + Inches(0.16)

    # ── Bullets ───────────────────────────────────────────────────────────────
    # Fewer bullets → bigger font; many bullets → two-col with smaller font
    bullet_font = 13 if two_col else 17
    line_h      = Inches(0.48) if two_col else Inches(0.62)

    if two_col:
        half  = len(bullets) // 2 + len(bullets) % 2
        col_w = (W - Inches(1.1)) / 2
        for col_idx, col_bullets in enumerate([bullets[:half], bullets[half:]]):
            x = Inches(0.45) + col_idx * (col_w + Inches(0.2))
            y = y_start
            for b in col_bullets:
                indent = b.startswith("  –")
                _pptx_textbox(slide,
                              x + (Inches(0.3) if indent else 0),
                              y,
                              col_w - (Inches(0.3) if indent else 0),
                              line_h,
                              ("  – " if indent else "• ") + b.lstrip(),
                              bullet_font - (2 if indent else 0),
                              bold=not indent,
                              color=PT_BODY if indent else PT_NAVY)
                y += line_h
    else:
        y = y_start
        for b in bullets:
            indent = b.startswith("  –")
            _pptx_textbox(slide,
                          Inches(0.45) + (Inches(0.4) if indent else 0),
                          y,
                          W - Inches(0.9 if not indent else 1.3),
                          line_h,
                          ("  – " if indent else "• ") + b.lstrip(),
                          bullet_font - (3 if indent else 0),
                          bold=not indent,
                          color=PT_BODY if indent else PT_NAVY)
            y += line_h

    if page_num:
        _pptx_footer(slide, page_num)


def _pptx_team_slide(prs, team_members, page_num, logo_path=None):
    slide = _pptx_blank_slide(prs)
    W, H = PPTX_W, PPTX_H
    _pptx_bg(slide, PT_WHITE)

    # Large navy title
    _pptx_textbox(slide, Inches(0.45), Inches(0.22),
                  W - Inches(0.9), Inches(0.92),
                  "Our Team of Experts", 32, bold=True, color=PT_NAVY)
    # Thin green divider
    _pptx_rect(slide, Inches(0.45), Inches(1.18),
               W - Inches(0.9), Inches(0.04), PT_GREEN)

    content_h = H - Inches(1.30) - Inches(0.58)
    row_h = content_h / len(team_members)

    for i, (name, role, bio) in enumerate(team_members):
        y = Inches(1.30) + i * row_h

        # Alternating light row stripe
        if i % 2 == 0:
            _pptx_rect(slide, Inches(0.35), y + Inches(0.04),
                       W - Inches(0.7), row_h - Inches(0.06), PT_LGRAY)

        # Name
        _pptx_textbox(slide, Inches(0.5), y + Inches(0.08),
                      Inches(4.0), Inches(0.46),
                      name, 17, bold=True, color=PT_NAVY)
        # Role (green)
        _pptx_textbox(slide, Inches(0.5), y + Inches(0.52),
                      Inches(4.0), Inches(0.34),
                      role, 13, bold=False, color=PT_GREEN)
        # Bio
        _pptx_textbox(slide, Inches(4.8), y + Inches(0.08),
                      W - Inches(5.3), row_h - Inches(0.16),
                      bio, 13, color=PT_BODY, wrap=True)

    _pptx_footer(slide, page_num)


def _pptx_services_slide(prs, company, services, notes, page_num, logo_path=None):
    slide = _pptx_blank_slide(prs)
    W, H = PPTX_W, PPTX_H
    _pptx_bg(slide, PT_WHITE)

    # Title — consistent 32pt (matches other content slides)
    _pptx_textbox(slide, Inches(0.45), Inches(0.22),
                  W - Inches(0.9), Inches(0.85),
                  "Recommended Services", 32, bold=True, color=PT_NAVY)
    # Subtitle — 13pt
    _pptx_textbox(slide, Inches(0.45), Inches(1.10),
                  W - Inches(0.9), Inches(0.34),
                  f"Tailored engagement plan for {company}", 13, color=PT_BODY)
    # Green divider
    _pptx_rect(slide, Inches(0.45), Inches(1.48),
               W - Inches(0.9), Inches(0.04), PT_GREEN)

    content_top  = Inches(1.90)   # extra breathing room below the green divider
    footer_h     = Inches(0.65)   # space reserved for footer
    note_h       = Inches(0.55) if notes else Inches(0)
    note_gap     = Inches(0.10) if notes else Inches(0)
    content_bottom = H - footer_h - note_h - note_gap

    half        = (len(services) + 1) // 2
    col_w       = (W - Inches(1.4)) / 2
    # Scale line height so bullets never overflow the available space
    available_h = content_bottom - content_top
    line_h      = min(Inches(0.65), available_h / max(half, 1))
    line_h      = max(line_h, Inches(0.38))   # floor to keep text readable
    bullet_font = 20 if line_h >= Inches(0.54) else (17 if line_h >= Inches(0.44) else 14)

    for col_idx, col_svcs in enumerate([services[:half], services[half:]]):
        x = Inches(0.45) + col_idx * (col_w + Inches(0.5))
        y = content_top
        for svc in col_svcs:
            _pptx_textbox(slide, x, y, col_w, line_h,
                          f"• {svc}", bullet_font, bold=True, color=PT_NAVY)
            y += line_h

    if notes:
        note_y = H - footer_h - note_h
        _pptx_textbox(slide, Inches(0.45), note_y,
                      W - Inches(0.9), note_h,
                      notes, 14, italic=True, color=PT_BODY, wrap=True)

    _pptx_footer(slide, page_num)


def _pptx_investment_slide(prs, company, costs, page_num, logo_path=None):
    slide = _pptx_blank_slide(prs)
    W, H = PPTX_W, PPTX_H
    _pptx_bg(slide, PT_WHITE)

    # Title — consistent 32pt (matches other content slides)
    _pptx_textbox(slide, Inches(0.45), Inches(0.22),
                  W - Inches(0.9), Inches(0.85),
                  "Proposed Investment", 32, bold=True, color=PT_NAVY)
    # Subtitle — 13pt
    _pptx_textbox(slide, Inches(0.45), Inches(1.10),
                  W - Inches(0.9), Inches(0.34),
                  f"Fee structure for {company}", 13, color=PT_BODY)
    # Green divider
    _pptx_rect(slide, Inches(0.45), Inches(1.48),
               W - Inches(0.9), Inches(0.04), PT_GREEN)

    summary_keys  = ("total investment", "monthly retainer", "one-time fees")
    item_costs    = [l for l in costs if not any(l.lower().startswith(k) for k in summary_keys)]
    summary_lines = [l for l in costs if any(l.lower().startswith(k) for k in summary_keys)]
    total_lines   = [l for l in costs if l.lower().startswith("total")]
    subtotals     = [l for l in summary_lines if not l.lower().startswith("total")]

    # ── Split items into Monthly vs One-Time groups ───────────────────────────
    monthly_items = [l for l in item_costs if '/mo' in l]
    onetime_items = [l for l in item_costs if '/mo' not in l and 'one-time' in l.lower()]
    other_items   = [l for l in item_costs if '/mo' not in l and 'one-time' not in l.lower()]
    # If no explicit type detected, treat all as a single flat list
    show_groups   = bool(monthly_items and onetime_items)
    # Groups to render in order
    groups = []
    if show_groups:
        if monthly_items:
            groups.append(("Monthly Services", PT_NAVY,              monthly_items))
        if onetime_items:
            groups.append(("One-Time Services", RGBColor(0x0E, 0x68, 0x82), onetime_items))
        if other_items:
            groups.append(("Additional Services", PT_NAVY,            other_items))
    else:
        groups.append((None, None, item_costs))   # no headers

    # ── Dynamic row sizing — fit everything in available space ────────────────
    content_top    = Inches(1.62)
    disclaimer_h   = Inches(0.40)
    footer_h       = Inches(0.65)
    content_bottom = H - footer_h - disclaimer_h - Inches(0.08)
    available_h    = content_bottom - content_top

    n_section_hdrs = sum(1 for lbl, _, _ in groups if lbl)
    hdr_h_factor   = 0.45   # section headers are 45% of a normal row
    gap_sub        = Inches(0.06) if subtotals   else Inches(0)
    gap_total      = Inches(0.10) if total_lines else Inches(0)
    gap_groups     = Inches(0.06) * max(0, n_section_hdrs - 1)  # gap between groups

    n_units = (len(item_costs)
               + n_section_hdrs * hdr_h_factor
               + len(subtotals) * 0.8
               + (1.3 if total_lines else 0))
    if n_units > 0:
        base_row_h = (available_h - gap_sub - gap_total - gap_groups) / n_units
        base_row_h = min(base_row_h, Inches(0.72))
        base_row_h = max(base_row_h, Inches(0.36))
    else:
        base_row_h = Inches(0.65)

    section_hdr_h = base_row_h * hdr_h_factor
    sub_row_h     = base_row_h * 0.8
    total_row_h   = base_row_h * 1.3
    font_item     = max(12, min(24, round(base_row_h  / Inches(0.72) * 24)))
    font_hdr      = max(10, min(16, round(section_hdr_h / Inches(0.32) * 14)))
    font_sub      = max(11, min(20, round(sub_row_h   / Inches(0.54) * 20)))
    font_total    = max(14, min(26, round(total_row_h  / Inches(0.82) * 26)))

    y = content_top

    def _render_item_row(line, stripe_idx, row_h, fnt):
        nonlocal y
        label, raw_amount = (line.split(":", 1) if ":" in line else (line, ""))
        # Strip type suffix for clean display; keep /mo for monthly clarity
        display_amount = raw_amount.strip()
        display_amount = display_amount.replace(' one-time', '').replace('one-time', '')
        # No alternating stripe — clean white background for line items
        padding = min(Inches(0.12), row_h * 0.18)
        _pptx_textbox(slide, Inches(0.65), y + padding,
                      Inches(7.5), row_h - padding,
                      label.strip(), fnt, bold=True, color=PT_NAVY)
        _pptx_textbox(slide, W - Inches(4.2), y + padding,
                      Inches(3.7), row_h - padding,
                      display_amount, fnt, bold=True,
                      color=RGBColor(0x3A, 0x7A, 0x3A), align=PP_ALIGN.RIGHT)
        y += row_h

    first_group = True
    for group_label, hdr_color, group_items in groups:
        if not group_items:
            continue
        if not first_group:
            y += Inches(0.06)   # small gap between groups
        first_group = False
        # Section header row
        if group_label:
            _pptx_rect(slide, Inches(0.4), y, W - Inches(0.8), section_hdr_h, hdr_color)
            padding = min(Inches(0.06), section_hdr_h * 0.2)
            hdr_box = _pptx_textbox(slide, Inches(0.65), y + padding,
                                    W - Inches(1.3), section_hdr_h,
                                    group_label.upper(), font_hdr,
                                    bold=True, color=PT_WHITE)
            hdr_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            y += section_hdr_h
        # Item rows
        for stripe_i, line in enumerate(group_items):
            _render_item_row(line, stripe_i, base_row_h, font_item)

    # ── Subtotal rows (Monthly Retainer / One-Time Fees) ─────────────────────
    if subtotals:
        y += gap_sub
        for line in subtotals:
            label, amount = (line.split(":", 1) if ":" in line else (line, ""))
            _pptx_rect(slide, Inches(0.4), y, W - Inches(0.8), sub_row_h, RGBColor(0xE8, 0xF0, 0xF8))
            padding = min(Inches(0.08), sub_row_h * 0.15)
            sub_lbl = _pptx_textbox(slide, Inches(0.65), y + padding,
                                    Inches(7.5), sub_row_h,
                                    label.strip(), font_sub, bold=True, italic=True, color=PT_NAVY)
            sub_lbl.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            sub_amt = _pptx_textbox(slide, W - Inches(4.2), y + padding,
                                    Inches(3.7), sub_row_h,
                                    "  " + amount.strip(), font_sub, bold=True,
                                    color=RGBColor(0x0E, 0x68, 0x82), align=PP_ALIGN.RIGHT)
            sub_amt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            y += sub_row_h

    # ── Total row ─────────────────────────────────────────────────────────────
    if total_lines:
        label, amount = (total_lines[0].split(":", 1) if ":" in total_lines[0] else (total_lines[0], ""))
        y += gap_total
        _pptx_rect(slide, Inches(0.4), y, W - Inches(0.8), total_row_h, PT_NAVY)
        padding = min(Inches(0.14), total_row_h * 0.18)
        tot_lbl = _pptx_textbox(slide, Inches(0.65), y + padding,
                                Inches(7.5), total_row_h,
                                label.strip(), font_total, bold=True, color=PT_WHITE)
        tot_lbl.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tot_amt = _pptx_textbox(slide, W - Inches(4.2), y + padding,
                                Inches(3.7), total_row_h,
                                "  " + amount.strip(), font_total, bold=True,
                                color=PT_GREEN, align=PP_ALIGN.RIGHT)
        tot_amt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Disclaimer — anchored just above the footer, never overlaps content
    _pptx_textbox(slide, Inches(0.4), H - footer_h - disclaimer_h,
                  W - Inches(0.8), disclaimer_h,
                  "All fees are subject to final scope confirmation. "
                  "Travel expenses billed at cost.",
                  12, italic=True, color=PT_MGRAY)

    _pptx_footer(slide, page_num)


def _pptx_image_slide(prs, bg_image_path):
    """Create a slide using a full-bleed background image (matches PDF template exactly)."""
    slide = _pptx_blank_slide(prs)
    if bg_image_path and os.path.exists(bg_image_path):
        slide.shapes.add_picture(bg_image_path, 0, 0, PPTX_W, PPTX_H)
    return slide


def build_proposal_pptx(company, contact, date,
                        services, costs, notes,
                        output_path, logo_path):
    prs = Presentation()
    prs.slide_width  = PPTX_W
    prs.slide_height = PPTX_H

    # Resolve assets directory (same folder as the logo)
    if logo_path and os.path.exists(logo_path):
        assets_dir = os.path.dirname(logo_path)
    else:
        assets_dir = os.path.join(os.path.dirname(__file__), '..', 'assets')

    # 1. Cover slide (dynamic — company/contact/date change per proposal)
    _pptx_cover(prs, company, contact, date, logo_path)

    # 2. Standard content slides — use pre-rendered PNG backgrounds from the
    #    PDF template so the PPTX looks pixel-perfect identical to the PDF.
    #    Template pages 2–8 map to slide_bg_2.png … slide_bg_8.png.
    for page_num in range(2, 9):
        bg_path = os.path.join(assets_dir, f'slide_bg_{page_num}.png')
        if os.path.exists(bg_path):
            _pptx_image_slide(prs, bg_path)
        else:
            # Fallback: text-based slide if background images are missing
            slide_idx = page_num - 2          # 0-based into STANDARD_SLIDES
            if slide_idx < len(STANDARD_SLIDES):
                slide_def = STANDARD_SLIDES[slide_idx]
                if "team" in slide_def:
                    _pptx_team_slide(prs, slide_def["team"], page_num,
                                     logo_path=logo_path)
                else:
                    two_col = len(slide_def.get("bullets", [])) > 6
                    _pptx_content_slide(
                        prs,
                        title    = slide_def["title"],
                        bullets  = slide_def.get("bullets", []),
                        subtitle = slide_def.get("subtitle"),
                        page_num = page_num,
                        two_col  = two_col,
                        logo_path= logo_path,
                    )

    # 3. Custom services slide (always dynamic — different per proposal)
    page = 9
    if services:
        _pptx_services_slide(prs, company, services, notes, page,
                             logo_path=logo_path)
        page += 1

    # 4. Investment slide (always dynamic)
    if costs:
        _pptx_investment_slide(prs, company, costs, page, logo_path=logo_path)

    prs.save(output_path)
    print(f"PPTX saved → {output_path}")


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="Generate Xcelerate proposal — outputs both PDF and PPTX")
    parser.add_argument("--company",  default="[COMPANY NAME]")
    parser.add_argument("--contact",  default="")
    parser.add_argument("--date",     default="[DATE]")
    parser.add_argument("--base-pdf", default="",
                        help="Path to the base proposal template PDF")
    parser.add_argument("--services", default="",
                        help="Pipe-separated services: 'Service A|Service B'")
    parser.add_argument("--costs",    default="",
                        help="Pipe-separated cost lines: 'Label: $amount|...'")
    parser.add_argument("--notes",    default="",
                        help="Optional custom paragraph for the services slide")
    parser.add_argument("--output",   default="Xcelerate_Proposal.pdf",
                        help="Output PDF path (.pptx is saved alongside automatically)")
    args = parser.parse_args()

    script_dir = os.path.dirname(os.path.abspath(__file__))
    logo       = _logo_path(script_dir)

    services = [s.strip() for s in args.services.split("|") if s.strip()] if args.services else []
    costs    = [c.strip() for c in args.costs.split("|")    if c.strip()] if args.costs    else []

    # ── PDF ──────────────────────────────────────────────────────────────────
    build_proposal_pdf(
        company       = args.company,
        contact       = args.contact,
        date          = args.date,
        base_pdf_path = args.base_pdf,
        services      = services,
        costs         = costs,
        notes         = args.notes,
        output_path   = args.output,
        logo_path     = logo,
    )

    # ── PPTX ─────────────────────────────────────────────────────────────────
    pptx_out = _pptx_output_path(args.output)
    build_proposal_pptx(
        company     = args.company,
        contact     = args.contact,
        date        = args.date,
        services    = services,
        costs       = costs,
        notes       = args.notes,
        output_path = pptx_out,
        logo_path   = logo,
    )


if __name__ == "__main__":
    main()
